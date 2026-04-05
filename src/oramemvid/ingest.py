import hashlib
import os
import oracledb
from oramemvid.embeddings import EmbeddingProvider
from oramemvid.frames import create_frame
from oramemvid.llm import LLMProvider
from oramemvid.memory_cards import create_memory_card


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == ".pdf":
        import pymupdf
        doc = pymupdf.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text
    elif ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)
    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    texts.append(" | ".join(cells))
        wb.close()
        return "\n".join(texts)
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    else:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()


def chunk_text(text: str, chunk_size: int = 512, chunk_overlap: int = 50) -> list[str]:
    words = text.split()
    if len(words) <= chunk_size:
        return [text.strip()] if text.strip() else []
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start = end - chunk_overlap
    return chunks


def _hash_file(file_path: str) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()


def ingest_file(
    conn: oracledb.Connection, file_path: str, provider: EmbeddingProvider,
    llm: LLMProvider | None = None, chunk_size: int = 512, chunk_overlap: int = 50,
) -> dict:
    filename = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lstrip(".").lower()
    file_hash = _hash_file(file_path)
    cursor = conn.cursor()

    cursor.execute("SELECT doc_id, total_frames FROM documents WHERE file_hash = :hash", {"hash": file_hash})
    existing = cursor.fetchone()
    if existing:
        return {"doc_id": existing[0], "filename": filename, "total_frames": existing[1], "skipped": True}

    doc_id_var = cursor.var(oracledb.NUMBER)
    cursor.execute("""
        INSERT INTO documents (filename, doc_type, file_hash)
        VALUES (:filename, :doc_type, :hash) RETURNING doc_id INTO :doc_id
    """, {"filename": filename, "doc_type": ext or "txt", "hash": file_hash, "doc_id": doc_id_var})
    conn.commit()
    doc_id = int(doc_id_var.getvalue()[0])

    text = extract_text(file_path)
    chunks = chunk_text(text, chunk_size, chunk_overlap)

    frame_ids = []
    for i, chunk in enumerate(chunks):
        uri = f"file://{filename}/chunk/{i}"
        frame_id = create_frame(
            conn=conn, uri=uri, content=chunk, provider=provider,
            title=f"{filename} (chunk {i})", doc_id=doc_id,
        )
        frame_ids.append(frame_id)
        if llm is not None:
            cards = llm.extract_memories(chunk)
            for card in cards:
                create_memory_card(
                    conn=conn, entity=card.get("entity", "unknown"),
                    slot=card.get("slot", "unknown"), value=card.get("value", ""),
                    kind=card.get("kind", "Fact"), source_frame_id=frame_id,
                    confidence=card.get("confidence", 1.0),
                )

    cursor.execute("UPDATE documents SET total_frames = :count WHERE doc_id = :id", {"count": len(frame_ids), "id": doc_id})
    conn.commit()
    return {"doc_id": doc_id, "filename": filename, "total_frames": len(frame_ids), "frame_ids": frame_ids, "skipped": False}


def ingest_text(
    conn: oracledb.Connection, text: str, uri: str, provider: EmbeddingProvider,
    llm: LLMProvider | None = None, title: str | None = None,
    chunk_size: int = 512, chunk_overlap: int = 50,
) -> dict:
    chunks = chunk_text(text, chunk_size, chunk_overlap)
    frame_ids = []
    for i, chunk in enumerate(chunks):
        chunk_uri = f"{uri}/chunk/{i}" if len(chunks) > 1 else uri
        frame_id = create_frame(
            conn=conn, uri=chunk_uri, content=chunk, provider=provider,
            title=title or f"text (chunk {i})",
        )
        frame_ids.append(frame_id)
        if llm is not None:
            cards = llm.extract_memories(chunk)
            for card in cards:
                create_memory_card(
                    conn=conn, entity=card.get("entity", "unknown"),
                    slot=card.get("slot", "unknown"), value=card.get("value", ""),
                    kind=card.get("kind", "Fact"), source_frame_id=frame_id,
                    confidence=card.get("confidence", 1.0),
                )
    return {"uri": uri, "total_frames": len(frame_ids), "frame_ids": frame_ids}
